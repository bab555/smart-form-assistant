"""
文档处理服务 - 支持 Excel/Word/PPT/PDF 解析
"""
import io
import base64
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import pandas as pd
from app.core.logger import app_logger as logger
from app.core.config import settings
from app.services.aliyun_llm import llm_service


class DocumentService:
    """文档处理服务"""
    
    # 支持的文件类型
    SUPPORTED_EXTENSIONS = {
        'excel': ['.xlsx', '.xls', '.csv'],
        'word': ['.docx', '.doc'],
        'ppt': ['.pptx', '.ppt'],
        'pdf': ['.pdf'],
        'image': ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp']
    }
    
    def __init__(self):
        logger.info("文档处理服务初始化完成")
    
    def get_file_type(self, filename: str) -> Optional[str]:
        """
        获取文件类型
        
        Args:
            filename: 文件名
            
        Returns:
            文件类型: excel/word/ppt/pdf/image 或 None
        """
        ext = Path(filename).suffix.lower()
        for file_type, extensions in self.SUPPORTED_EXTENSIONS.items():
            if ext in extensions:
                return file_type
        return None
    
    def is_supported(self, filename: str) -> bool:
        """检查文件是否支持"""
        return self.get_file_type(filename) is not None
    
    async def extract_data(
        self,
        file_content: bytes,
        filename: str,
        template_columns: Optional[List[Dict]] = None
    ) -> Dict[str, Any]:
        """
        从文档中提取数据（统一入口）
        
        Args:
            file_content: 文件二进制内容
            filename: 文件名
            template_columns: 模板列定义（用于字段映射）
            
        Returns:
            统一的提取结果结构
        """
        file_type = self.get_file_type(filename)
        
        if not file_type:
            raise ValueError(f"不支持的文件类型: {filename}")
        
        logger.info(f"[文档处理] 开始处理: {filename} (类型: {file_type})")
        
        try:
            if file_type == 'excel':
                return await self._extract_from_excel(file_content, filename, template_columns)
            elif file_type in ['word', 'ppt', 'pdf']:
                return await self._extract_from_document(file_content, filename, file_type, template_columns)
            elif file_type == 'image':
                return await self._extract_from_image(file_content, filename, template_columns)
            else:
                raise ValueError(f"未实现的文件类型处理: {file_type}")
                
        except Exception as e:
            logger.error(f"[文档处理] 处理失败: {str(e)}")
            raise
    
    async def _extract_from_excel(
        self,
        file_content: bytes,
        filename: str,
        template_columns: Optional[List[Dict]] = None
    ) -> Dict[str, Any]:
        """
        从 Excel 提取数据（直接解析，保留结构）
        
        Excel 是结构化数据，直接读取比转图片更精准
        """
        logger.info(f"[Excel处理] 开始解析: {filename}")
        
        try:
            # 读取 Excel
            ext = Path(filename).suffix.lower()
            if ext == '.csv':
                df = pd.read_csv(io.BytesIO(file_content))
            else:
                df = pd.read_excel(io.BytesIO(file_content), engine='openpyxl')
            
            # 清理数据
            df = df.dropna(how='all')  # 删除全空行
            df = df.fillna('')  # 空值填充为空字符串
            
            logger.info(f"[Excel处理] 读取成功: {len(df)} 行, {len(df.columns)} 列")
            logger.info(f"[Excel处理] 列名: {list(df.columns)}")
            
            # 如果有模板列定义，尝试智能映射
            if template_columns:
                rows = await self._map_excel_to_template(df, template_columns)
            else:
                # 无模板，直接转换
                rows = self._dataframe_to_rows(df)
            
            return {
                "success": True,
                "file_type": "excel",
                "rows": rows,
                "raw_columns": list(df.columns),
                "row_count": len(rows),
                "message": f"成功从 Excel 提取 {len(rows)} 行数据"
            }
            
        except Exception as e:
            logger.error(f"[Excel处理] 解析失败: {str(e)}")
            raise ValueError(f"Excel 解析失败: {str(e)}")
    
    def _dataframe_to_rows(self, df: pd.DataFrame) -> List[List[Dict]]:
        """
        将 DataFrame 转换为前端期望的 FormItem[][] 格式
        """
        rows = []
        columns = list(df.columns)
        
        for idx, row in df.iterrows():
            row_items = []
            for col in columns:
                value = row[col]
                # 转换为字符串，处理各种类型
                if pd.isna(value):
                    str_value = ""
                elif isinstance(value, (int, float)):
                    str_value = str(int(value)) if value == int(value) else str(value)
                else:
                    str_value = str(value)
                
                row_items.append({
                    "key": str(col).replace(' ', '_').lower(),
                    "label": str(col),
                    "value": str_value,
                    "original_text": str_value,
                    "confidence": 1.0,  # Excel 直接读取，置信度为1
                    "is_ambiguous": False,
                    "candidates": None,
                    "data_type": "number" if isinstance(value, (int, float)) else "string"
                })
            
            if row_items:
                rows.append(row_items)
        
        return rows
    
    async def _map_excel_to_template(
        self,
        df: pd.DataFrame,
        template_columns: List[Dict]
    ) -> List[List[Dict]]:
        """
        智能映射 Excel 列到模板字段（使用 LLM）
        """
        excel_columns = list(df.columns)
        template_keys = [col.get('key') for col in template_columns]
        template_labels = [col.get('label') for col in template_columns]
        
        # 构建映射提示
        prompt = f"""你是一个数据映射助手。请将 Excel 的列名映射到目标模板字段。

Excel 列名: {excel_columns}

目标模板字段:
{chr(10).join([f"- {col.get('key')}: {col.get('label')}" for col in template_columns])}

请返回 JSON 格式的映射关系，key 是 Excel 列名，value 是模板字段 key。
如果某列无法映射，value 设为 null。
只返回 JSON，不要其他内容。

示例格式: {{"Excel列名1": "template_key1", "Excel列名2": null}}"""

        try:
            response = await llm_service.call_calibration_model(prompt)
            
            import json
            import re
            json_match = re.search(r'\{[^}]+\}', response, re.DOTALL)
            if json_match:
                mapping = json.loads(json_match.group())
                logger.info(f"[Excel处理] 列映射结果: {mapping}")
            else:
                # 映射失败，使用默认映射
                mapping = {}
                
        except Exception as e:
            logger.warning(f"[Excel处理] LLM 映射失败，使用直接转换: {e}")
            mapping = {}
        
        # 应用映射并转换数据
        rows = []
        for idx, row in df.iterrows():
            row_items = []
            
            for excel_col in excel_columns:
                value = row[excel_col]
                if pd.isna(value):
                    str_value = ""
                elif isinstance(value, (int, float)):
                    str_value = str(int(value)) if value == int(value) else str(value)
                else:
                    str_value = str(value)
                
                # 查找映射的 key
                mapped_key = mapping.get(str(excel_col))
                if mapped_key and mapped_key in template_keys:
                    key = mapped_key
                    # 找到对应的 label
                    label = next((col.get('label') for col in template_columns if col.get('key') == key), str(excel_col))
                else:
                    key = str(excel_col).replace(' ', '_').lower()
                    label = str(excel_col)
                
                row_items.append({
                    "key": key,
                    "label": label,
                    "value": str_value,
                    "original_text": str_value,
                    "confidence": 1.0,
                    "is_ambiguous": False,
                    "candidates": None,
                    "data_type": "number" if isinstance(value, (int, float)) else "string"
                })
            
            if row_items:
                rows.append(row_items)
        
        return rows
    
    async def _extract_from_document(
        self,
        file_content: bytes,
        filename: str,
        file_type: str,
        template_columns: Optional[List[Dict]] = None
    ) -> Dict[str, Any]:
        """
        从 Word/PPT/PDF 提取数据（转图片 -> VL 识别）
        """
        logger.info(f"[文档处理] 开始转换为图片: {filename}")
        
        try:
            # 转换为图片
            images = await self._convert_to_images(file_content, filename, file_type)
            
            if not images:
                raise ValueError("文档转图片失败，未生成任何图片")
            
            logger.info(f"[文档处理] 转换完成，共 {len(images)} 页")
            
            # 使用 Qwen-VL 识别（只处理第一页，避免 token 过多）
            first_image = images[0]
            
            # 构建识别提示
            prompt = """请识别这张图片中的表格或数据内容。
如果是表格，请提取每一行的数据，返回 JSON 格式：
{"rows": [{"字段1": "值1", "字段2": "值2"}, ...]}

如果是普通文档，请提取关键信息，返回 JSON 格式：
{"rows": [{"内容": "提取的文本内容"}]}

只返回 JSON，不要其他内容。"""

            # 将图片转为 base64
            image_base64 = base64.b64encode(first_image).decode('utf-8')
            image_url = f"data:image/png;base64,{image_base64}"
            
            # 调用 VL 模型
            result_text = await llm_service.call_vl_model(image_url, prompt)
            
            # 解析结果
            import json
            import re
            json_match = re.search(r'\{[\s\S]*\}', result_text)
            if json_match:
                extracted = json.loads(json_match.group())
                raw_rows = extracted.get('rows', [])
            else:
                # 解析失败，返回原始文本
                raw_rows = [{"content": result_text}]
            
            # 转换为标准格式
            rows = []
            for row_data in raw_rows:
                row_items = []
                for key, value in row_data.items():
                    row_items.append({
                        "key": key.replace(' ', '_').lower(),
                        "label": key,
                        "value": str(value),
                        "original_text": str(value),
                        "confidence": 0.85,  # VL 识别置信度稍低
                        "is_ambiguous": False,
                        "candidates": None,
                        "data_type": "string"
                    })
                if row_items:
                    rows.append(row_items)
            
            return {
                "success": True,
                "file_type": file_type,
                "rows": rows,
                "page_count": len(images),
                "row_count": len(rows),
                "message": f"成功从 {file_type.upper()} 提取 {len(rows)} 行数据（共 {len(images)} 页）"
            }
            
        except Exception as e:
            logger.error(f"[文档处理] 处理失败: {str(e)}")
            raise
    
    async def convert_to_images(
        self,
        file_content: bytes,
        filename: str,
        file_type: str = None
    ) -> List[bytes]:
        """
        将文档转换为图片列表（公开接口）
        """
        if not file_type:
            file_type = self.get_file_type(filename)
            
        return await self._convert_to_images(file_content, filename, file_type)

    async def _convert_to_images(
        self,
        file_content: bytes,
        filename: str,
        file_type: str
    ) -> List[bytes]:
        """
        将文档转换为图片列表
        """
        images = []
        
        try:
            if file_type == 'pdf':
                images = self._pdf_to_images(file_content)
            elif file_type == 'word':
                images = await self._word_to_images(file_content, filename)
            elif file_type == 'ppt':
                images = await self._ppt_to_images(file_content, filename)
            
            return images
            
        except ImportError as e:
            logger.warning(f"[文档转换] 依赖缺失: {e}，尝试备用方案")
            # 如果转换库不可用，尝试直接提取文本
            return await self._fallback_extract_text(file_content, filename, file_type)
        except Exception as e:
            logger.error(f"[文档转换] 转换失败: {str(e)}")
            raise
    
    def _pdf_to_images(self, file_content: bytes) -> List[bytes]:
        """PDF 转图片"""
        try:
            from pdf2image import convert_from_bytes
            from PIL import Image
            
            # 转换 PDF 为图片列表
            pil_images = convert_from_bytes(file_content, dpi=150, first_page=1, last_page=3)
            
            images = []
            for img in pil_images:
                img_bytes = io.BytesIO()
                img.save(img_bytes, format='PNG')
                images.append(img_bytes.getvalue())
            
            return images
            
        except Exception as e:
            logger.error(f"[PDF转换] 失败: {e}")
            raise
    
    async def _word_to_images(self, file_content: bytes, filename: str) -> List[bytes]:
        """
        Word 转图片
        备用方案：提取文本内容，让 LLM 处理
        """
        try:
            from docx import Document
            from PIL import Image, ImageDraw, ImageFont
            
            # 读取 Word 文档
            doc = Document(io.BytesIO(file_content))
            
            # 提取所有段落文本
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())
            
            # 提取表格
            tables_text = []
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(row_text))
                tables_text.append("\n".join(table_rows))
            
            # 合并内容
            all_text = "\n".join(paragraphs)
            if tables_text:
                all_text += "\n\n[表格内容]\n" + "\n\n".join(tables_text)
            
            # 将文本渲染为图片（简单方案）
            img = self._text_to_image(all_text)
            
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG')
            
            return [img_bytes.getvalue()]
            
        except Exception as e:
            logger.error(f"[Word转换] 失败: {e}")
            # 返回空，后续走文本提取备用方案
            raise
    
    async def _ppt_to_images(self, file_content: bytes, filename: str) -> List[bytes]:
        """
        PPT 转图片
        备用方案：提取文本内容
        """
        try:
            from pptx import Presentation
            from PIL import Image, ImageDraw
            
            # 读取 PPT
            prs = Presentation(io.BytesIO(file_content))
            
            images = []
            for slide_idx, slide in enumerate(prs.slides[:3]):  # 只处理前3页
                # 提取幻灯片文本
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                
                slide_text = f"[幻灯片 {slide_idx + 1}]\n" + "\n".join(texts)
                
                # 渲染为图片
                img = self._text_to_image(slide_text)
                
                img_bytes = io.BytesIO()
                img.save(img_bytes, format='PNG')
                images.append(img_bytes.getvalue())
            
            return images
            
        except Exception as e:
            logger.error(f"[PPT转换] 失败: {e}")
            raise
    
    def _text_to_image(self, text: str, width: int = 800, padding: int = 20) -> 'Image':
        """
        将文本渲染为图片（用于 VL 模型识别）
        """
        from PIL import Image, ImageDraw, ImageFont
        
        # 限制文本长度
        if len(text) > 2000:
            text = text[:2000] + "\n...(内容过长，已截断)"
        
        # 计算图片高度
        lines = text.split('\n')
        line_height = 24
        height = max(200, len(lines) * line_height + padding * 2)
        
        # 创建图片
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)
        
        # 尝试使用系统字体
        try:
            font = ImageFont.truetype("arial.ttf", 16)
        except:
            font = ImageFont.load_default()
        
        # 绘制文本
        y = padding
        for line in lines:
            draw.text((padding, y), line, fill='black', font=font)
            y += line_height
        
        return img
    
    async def _fallback_extract_text(
        self,
        file_content: bytes,
        filename: str,
        file_type: str
    ) -> List[bytes]:
        """
        备用方案：直接提取文本，渲染为图片
        """
        text = ""
        
        try:
            if file_type == 'word':
                from docx import Document
                doc = Document(io.BytesIO(file_content))
                text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                
            elif file_type == 'ppt':
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_content))
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slides_text.append(shape.text)
                text = "\n".join(slides_text)
                
        except Exception as e:
            text = f"[文档解析失败: {str(e)}]"
        
        if not text:
            text = "[未能提取文档内容]"
        
        # 渲染为图片
        img = self._text_to_image(text)
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        
        return [img_bytes.getvalue()]
    
    async def _extract_from_image(
        self,
        file_content: bytes,
        filename: str,
        template_columns: Optional[List[Dict]] = None
    ) -> Dict[str, Any]:
        """
        从图片提取数据（直接调用 VL）
        """
        logger.info(f"[图片处理] 开始识别: {filename}")
        
        # 转为 base64
        image_base64 = base64.b64encode(file_content).decode('utf-8')
        
        # 判断图片格式
        ext = Path(filename).suffix.lower()
        mime_type = {
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.gif': 'image/gif',
            '.webp': 'image/webp',
            '.bmp': 'image/bmp'
        }.get(ext, 'image/png')
        
        image_url = f"data:{mime_type};base64,{image_base64}"
        
        # 构建提示
        prompt = """请识别这张图片中的表格或订单数据。
提取每一行的信息，返回 JSON 格式：
{"rows": [{"商品名称": "xxx", "数量": "xxx", "单位": "xxx", "单价": "xxx"}, ...]}

如果字段名不确定，使用你认为最合适的名称。
只返回 JSON，不要其他内容。"""

        result_text = await llm_service.call_vl_model(image_url, prompt)
        
        # 解析结果
        import json
        import re
        json_match = re.search(r'\{[\s\S]*\}', result_text)
        if json_match:
            extracted = json.loads(json_match.group())
            raw_rows = extracted.get('rows', [])
        else:
            raw_rows = []
        
        # 转换为标准格式
        rows = []
        for row_data in raw_rows:
            row_items = []
            for key, value in row_data.items():
                row_items.append({
                    "key": key.replace(' ', '_').lower(),
                    "label": key,
                    "value": str(value),
                    "original_text": str(value),
                    "confidence": 0.9,
                    "is_ambiguous": False,
                    "candidates": None,
                    "data_type": "string"
                })
            if row_items:
                rows.append(row_items)
        
        return {
            "success": True,
            "file_type": "image",
            "rows": rows,
            "row_count": len(rows),
            "message": f"成功从图片提取 {len(rows)} 行数据"
        }


# 全局单例
document_service = DocumentService()

